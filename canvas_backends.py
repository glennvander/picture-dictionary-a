"""
Two drawing backends behind one small interface, so the book layout is written
once and emitted identically as an editable .pptx and as a print-ready .pdf.

Coordinates are in inches, origin top-left, y increasing downward (PowerPoint
convention). The PDF backend flips y internally.

Interface:
    page(w, h)                      start a new page
    rect(x, y, w, h, color)
    rounded(x, y, w, h, color, r)
    image(path, x, y, w, h)         centre-crop to fill the box
    text(x, y, w, s, size, color, bold, italic, align, font, leading)
    save(path)

`color` is a "#rrggbb" string. `font` is "display" or "body".
"""

import os
from PIL import Image

DISPLAY = "display"
BODY = "body"


def _hex(c):
    c = c.lstrip("#")
    return tuple(int(c[i:i + 2], 16) for i in (0, 2, 4))


def _crop_box(path, bw, bh):
    """Return PIL box that centre-crops the image to the box aspect ratio."""
    with Image.open(path) as im:
        iw, ih = im.size
    box_ar, img_ar = bw / bh, iw / ih
    if img_ar > box_ar:
        nw = int(ih * box_ar)
        off = (iw - nw) // 2
        return (off, 0, off + nw, ih)
    nh = int(iw / box_ar)
    off = (ih - nh) // 2
    return (0, off, iw, off + nh)


# --------------------------------------------------------------------- PPTX

class PptxCanvas:
    FONTS = {DISPLAY: "Trebuchet MS", BODY: "Verdana"}

    def __init__(self):
        from pptx import Presentation
        self.prs = Presentation()
        self.slide = None
        self._sized = False

    def page(self, w, h):
        from pptx.util import Inches
        if not self._sized:
            self.prs.slide_width = Inches(w)
            self.prs.slide_height = Inches(h)
            self._sized = True
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _shape(self, kind, x, y, w, h, color):
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
        s = self.slide.shapes.add_shape(kind, Inches(x), Inches(y),
                                        Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(*_hex(color))
        s.line.fill.background()
        s.shadow.inherit = False
        return s

    def rect(self, x, y, w, h, color):
        from pptx.enum.shapes import MSO_SHAPE
        self._shape(MSO_SHAPE.RECTANGLE, x, y, w, h, color)

    def rounded(self, x, y, w, h, color, r=0.12):
        from pptx.enum.shapes import MSO_SHAPE
        s = self._shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, color)
        try:
            s.adjustments[0] = min(0.5, r / min(w, h))
        except Exception:
            pass

    def image(self, path, x, y, w, h):
        from pptx.util import Inches
        with Image.open(path) as im:
            iw, ih = im.size
        box_ar, img_ar = w / h, iw / ih
        pic = self.slide.shapes.add_picture(path, Inches(x), Inches(y),
                                            width=Inches(w), height=Inches(h))
        if img_ar > box_ar:
            f = (1 - box_ar / img_ar) / 2
            pic.crop_left = pic.crop_right = f
        elif img_ar < box_ar:
            f = (1 - img_ar / box_ar) / 2
            pic.crop_top = pic.crop_bottom = f

    def text(self, x, y, w, s, size, color, bold=False, italic=False,
             align="left", font=BODY, leading=1.2):
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        amap = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT}
        tb = self.slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w),
                                           Inches(size / 72 * leading * 1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        for i, line in enumerate(s.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = amap[align]
            p.line_spacing = leading
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = self.FONTS[font]
            r.font.color.rgb = RGBColor(*_hex(color))

    def save(self, path):
        os.makedirs(os.path.dirname(path), exist_ok=True)
        self.prs.save(path)


# ---------------------------------------------------------------------- PDF

class PdfCanvas:
    # Helvetica is metrically close to the Office fonts and is a core PDF
    # font, so the PDF needs no embedded font files and opens anywhere.
    FONTS = {
        (DISPLAY, False, False): "Helvetica",
        (DISPLAY, True, False): "Helvetica-Bold",
        (DISPLAY, False, True): "Helvetica-Oblique",
        (DISPLAY, True, True): "Helvetica-BoldOblique",
        (BODY, False, False): "Helvetica",
        (BODY, True, False): "Helvetica-Bold",
        (BODY, False, True): "Helvetica-Oblique",
        (BODY, True, True): "Helvetica-BoldOblique",
    }

    def __init__(self):
        self.c = None
        self.h = 0
        self._pending = None

    def page(self, w, h):
        from reportlab.pdfgen import canvas as rc
        if self.c is None:
            self._pending = (w * 72, h * 72)
            self.c = rc.Canvas("/dev/null", pagesize=self._pending)
            self.c.setPageCompression(1)
        else:
            self.c.showPage()
        self.c.setPageSize((w * 72, h * 72))
        self.h = h

    def _y(self, y):
        return (self.h - y) * 72

    def _fill(self, color):
        r, g, b = _hex(color)
        self.c.setFillColorRGB(r / 255, g / 255, b / 255)

    def rect(self, x, y, w, h, color):
        self._fill(color)
        self.c.rect(x * 72, self._y(y + h), w * 72, h * 72, stroke=0, fill=1)

    def rounded(self, x, y, w, h, color, r=0.12):
        self._fill(color)
        self.c.roundRect(x * 72, self._y(y + h), w * 72, h * 72, r * 72,
                         stroke=0, fill=1)

    def image(self, path, x, y, w, h):
        # Write the cropped result as a JPEG on disk and hand reportlab the
        # path: given a .jpg file it embeds the JPEG stream directly instead of
        # re-encoding to raw RGB, which is the difference between a ~3MB and a
        # ~100MB book.
        import hashlib
        key = hashlib.md5(f"{path}|{w:.3f}|{h:.3f}".encode()).hexdigest()[:16]
        cache = os.path.join(os.path.dirname(path), ".pdfcrop")
        os.makedirs(cache, exist_ok=True)
        dst = os.path.join(cache, key + ".jpg")
        if not os.path.exists(dst):
            with Image.open(path) as im:
                im.convert("RGB").crop(_crop_box(path, w, h)).save(
                    dst, "JPEG", quality=90, optimize=True)
        self.c.drawImage(dst, x * 72, self._y(y + h), w * 72, h * 72)

    def _wrap(self, s, font, size, w):
        from reportlab.pdfbase.pdfmetrics import stringWidth
        out = []
        for para in s.split("\n"):
            line = ""
            for word in para.split(" "):
                trial = (line + " " + word).strip()
                if stringWidth(trial, font, size) <= w * 72 or not line:
                    line = trial
                else:
                    out.append(line)
                    line = word
            out.append(line)
        return out

    def text(self, x, y, w, s, size, color, bold=False, italic=False,
             align="left", font=BODY, leading=1.2):
        f = self.FONTS[(font, bold, italic)]
        self._fill(color)
        self.c.setFont(f, size)
        lead = size * leading
        # first baseline sits one ascent below the top of the text box
        yy = self._y(y) - size * 0.83
        for line in self._wrap(s, f, size, w):
            if align == "center":
                self.c.drawCentredString(x * 72 + w * 72 / 2, yy, line)
            elif align == "right":
                self.c.drawRightString(x * 72 + w * 72, yy, line)
            else:
                self.c.drawString(x * 72, yy, line)
            yy -= lead

    def save(self, path):
        os.makedirs(os.path.dirname(path), exist_ok=True)
        self.c._filename = path
        self.c._doc.filename = path
        self.c.save()
